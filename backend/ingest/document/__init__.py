"""
document — 다양한 office/text 포맷 텍스트 추출 통합 모듈.

CLAUDE.md §2 (Raw-first / Loss-less): 추출 텍스트는 분석/검색용 derived data.
**원본 파일은 항상 별도 attachment 로 보존** (caller 책임 — backend/storage/local).

지원 포맷 (Phase 2.5):
- PDF  → backend.ingest.pdf._extract_pdf_text 재사용
- DOCX → python-docx
- PPTX → python-pptx
- TXT  → utf-8 / charset-normalizer 자동 감지
- MD   → utf-8

지원 안 함 (Phase 3+):
- DOC, PPT (구 binary 형식) — LibreOffice CLI 필요
- HWP/HWPX — pyhwp (한국 정부 문서)
- 이미지 — OCR (EasyOCR/Tesseract)

설계 — 진입점 2개:
- `guess_format(filename, mime_type)` → str — 포맷 식별
- `extract_text_from_bytes(data, filename, mime_type)` → DocumentExtract | None
  — 지원 안 하는 포맷이면 None (caller 는 attachment 만 저장)
"""

from __future__ import annotations

import io
import logging
import mimetypes
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker

from backend.db.connection import get_engine
from backend.db.repository import (
    find_item_by_hash,
    insert_attachment,
    insert_item,
)

# PDF 재사용 — backend/ingest/pdf 의 검증된 함수.
from backend.ingest.pdf import (
    _extract_pdf_text,
    _sanitize_text,
    _save_pdf_figures,
)
from backend.utils.external_ids import extract_external_ids
from backend.utils.hashing import sha256_text

logger = logging.getLogger(__name__)


@dataclass
class DocumentExtract:
    """포맷-agnostic 텍스트 추출 결과.

    `raw_text` 는 items.raw_content 로 들어감 — _sanitize_text 로 NUL byte 등 정제 완료.
    `title` 은 items.title 의 fallback (없으면 파일명).
    """
    raw_text: str
    title: str | None = None
    source_format: str = "unknown"        # "pdf" | "docx" | "pptx" | "txt" | "markdown"
    extractor: str = "unknown"            # "pypdf" | "pymupdf" | "python-docx" | "python-pptx" | "utf-8" | ...
    meta: dict[str, Any] = field(default_factory=dict)


# ──────────────────────────────────────────────────────────────
# Format 식별
# ──────────────────────────────────────────────────────────────


_EXT_TO_FORMAT: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".md": "markdown",
    ".markdown": "markdown",
    ".rst": "txt",          # rst 도 plain text 로 일단 (Phase 3+ 에 별 파서)
}


_MIME_TO_FORMAT: dict[str, str] = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/plain": "txt",
    "text/markdown": "markdown",
    "text/x-markdown": "markdown",
}


# Phase 3+ 에 추가 예정 — 지금은 attachment 만 저장
_KNOWN_UNSUPPORTED_FORMATS: set[str] = {
    "doc", "ppt", "hwp", "hwpx", "rtf",
    "image",  # png/jpg/gif/webp — caller 가 'image' 로 분류
}


def guess_format(
    filename: str | None = None, mime_type: str | None = None,
) -> str:
    """mime_type → 확장자 → "unknown" 순으로 포맷 식별.

    이미지 (mime_type 이 image/* 로 시작) 는 "image" 로 분류 — extract 는 None.
    """
    if mime_type:
        mt = mime_type.lower().split(";")[0].strip()
        if mt in _MIME_TO_FORMAT:
            return _MIME_TO_FORMAT[mt]
        if mt.startswith("image/"):
            return "image"

    if filename:
        ext = Path(filename).suffix.lower()
        if ext in _EXT_TO_FORMAT:
            return _EXT_TO_FORMAT[ext]
        # 구 binary office 명시적 표시 (Phase 3+ 에서 LibreOffice 로 처리)
        if ext == ".doc":
            return "doc"
        if ext == ".ppt":
            return "ppt"
        if ext in {".hwp", ".hwpx"}:
            return "hwp"
        if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}:
            return "image"

    return "unknown"


def is_supported(fmt: str) -> bool:
    """guess_format 결과가 텍스트 추출 가능한지."""
    return fmt in {"pdf", "docx", "pptx", "txt", "markdown"}


# ──────────────────────────────────────────────────────────────
# 포맷별 추출
# ──────────────────────────────────────────────────────────────


def _extract_docx(data: bytes) -> tuple[str, str | None, dict[str, Any]]:
    """python-docx — paragraph 텍스트 join."""
    from docx import Document   # type: ignore[import-not-found]

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text]
    # 표 안의 텍스트도 가져옴 (실무 docx 는 표가 많음)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                txt = cell.text
                if txt:
                    paragraphs.append(txt)
    body = "\n\n".join(paragraphs)

    # core_properties — title 우선, 없으면 None (caller 가 filename 으로 fallback)
    title = None
    try:
        cp = doc.core_properties
        if cp.title:
            title = cp.title.strip() or None
    except Exception:  # noqa: BLE001
        pass

    meta: dict[str, Any] = {
        "paragraph_count": len(paragraphs),
        "table_count": len(doc.tables),
    }
    return body, title, meta


def _extract_pptx(data: bytes) -> tuple[str, str | None, dict[str, Any]]:
    """python-pptx — 슬라이드별 텍스트 + notes 까지."""
    from pptx import Presentation   # type: ignore[import-not-found]

    pres = Presentation(io.BytesIO(data))
    slides_text: list[str] = []
    for i, slide in enumerate(pres.slides, start=1):
        parts: list[str] = [f"=== Slide {i} ==="]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in p.runs).strip()
                    if txt:
                        parts.append(txt)
        # 발표자 노트 — 학습 데이터로 가치 큼
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Notes] {notes_text}")
        slides_text.append("\n".join(parts))
    body = "\n\n".join(slides_text)

    # title slide 의 첫 텍스트를 title 로
    title = None
    if pres.slides:
        first = pres.slides[0]
        for shape in first.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().split("\n")[0][:200]
                break

    meta: dict[str, Any] = {"slide_count": len(pres.slides)}
    return body, title, meta


def _extract_text_file(data: bytes, fmt: str) -> tuple[str, dict[str, Any]]:
    """TXT/MD — 인코딩 자동 감지.

    순서:
    1. utf-8 (가장 흔함, 거의 모든 모던 환경)
    2. cp949 / euc-kr (한국 윈도우 메모장 / 옛 한글 텍스트 — charset-normalizer
       는 짧은 텍스트로 big5/cp949 구분 못 함, 한국 사용자 환경 우선)
    3. charset-normalizer (longer text 의 일반 자동 감지)
    4. 최후: utf-8 errors='replace' — loss-less 위반이지만 raw 는 attachment 로 보존됨
    """
    # 1. utf-8 — 가장 흔함
    try:
        return data.decode("utf-8"), {"encoding": "utf-8"}
    except UnicodeDecodeError:
        pass

    # 2. 한국어 인코딩 우선 시도 (사용자 한국 환경, 옛 메모장 파일)
    for enc in ("cp949", "euc-kr"):
        try:
            return data.decode(enc), {"encoding": enc}
        except UnicodeDecodeError:
            continue

    # 3. charset-normalizer 일반 자동 감지
    try:
        from charset_normalizer import from_bytes   # type: ignore[import-not-found]

        result = from_bytes(data).best()
        if result is not None:
            return str(result), {"encoding": result.encoding or "auto"}
    except Exception:  # noqa: BLE001
        pass

    # 4. 최후 fallback
    return data.decode("utf-8", errors="replace"), {"encoding": "utf-8-replace"}


# ──────────────────────────────────────────────────────────────
# 통합 진입점
# ──────────────────────────────────────────────────────────────


def extract_text_from_bytes(
    data: bytes,
    *,
    filename: str | None = None,
    mime_type: str | None = None,
) -> DocumentExtract | None:
    """파일 bytes → DocumentExtract.

    Args:
        data: 파일 raw bytes
        filename: 원 파일명 (확장자 + title fallback 용)
        mime_type: 알면 우선 사용 (없으면 filename 으로 추정)

    Returns:
        지원 포맷이면 DocumentExtract, 아니면 None (caller 가 attachment 만 저장).

    실패 시:
        파싱 예외 발생하면 logger.warning 후 None 반환 — caller 가 attachment 만
        저장하도록 graceful degrade.
    """
    if not data:
        return None

    # mime_type 모르면 filename 으로 추정
    if not mime_type and filename:
        guessed_mt, _ = mimetypes.guess_type(filename)
        mime_type = guessed_mt

    fmt = guess_format(filename, mime_type)

    if fmt == "pdf":
        try:
            text_out, info = _extract_pdf_text(data)
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "PDF 추출 실패 (file=%s, %s: %s)",
                filename, type(e).__name__, e,
            )
            return None
        if not text_out.strip():
            return None
        return DocumentExtract(
            raw_text=text_out,
            title=_filename_to_title(filename),
            source_format="pdf",
            extractor=info.get("extractor", "pdf"),
            meta=info,
        )

    if fmt == "docx":
        try:
            body, title, meta = _extract_docx(data)
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "DOCX 추출 실패 (file=%s, %s: %s)",
                filename, type(e).__name__, e,
            )
            return None
        if not body.strip():
            return None
        return DocumentExtract(
            raw_text=_sanitize_text(body),
            title=title or _filename_to_title(filename),
            source_format="docx",
            extractor="python-docx",
            meta=meta,
        )

    if fmt == "pptx":
        try:
            body, title, meta = _extract_pptx(data)
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "PPTX 추출 실패 (file=%s, %s: %s)",
                filename, type(e).__name__, e,
            )
            return None
        if not body.strip():
            return None
        return DocumentExtract(
            raw_text=_sanitize_text(body),
            title=title or _filename_to_title(filename),
            source_format="pptx",
            extractor="python-pptx",
            meta=meta,
        )

    if fmt in {"txt", "markdown"}:
        try:
            body, info = _extract_text_file(data, fmt)
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "TXT/MD 추출 실패 (file=%s, %s: %s)",
                filename, type(e).__name__, e,
            )
            return None
        if not body.strip():
            return None
        # MD 의 경우 첫 # heading 을 title 로 시도
        title = _filename_to_title(filename)
        if fmt == "markdown":
            for line in body.splitlines():
                line = line.strip()
                if line.startswith("# "):
                    title = line[2:].strip() or title
                    break
        return DocumentExtract(
            raw_text=_sanitize_text(body),
            title=title,
            source_format=fmt,
            extractor=info.get("encoding", "utf-8"),
            meta=info,
        )

    # unsupported (doc/ppt/hwp/image/unknown) — attachment 만 저장하도록 None
    return None


def _filename_to_title(filename: str | None) -> str | None:
    """확장자 제거 + 너무 길면 cut. 파일명만으로 부족하면 None."""
    if not filename:
        return None
    name = Path(filename).stem.strip()
    if not name:
        return None
    return name[:200]


# ──────────────────────────────────────────────────────────────
# ingest_document — 통합 진입점 (텔레그램 첨부, 향후 web upload, Slack 첨부 등)
# ──────────────────────────────────────────────────────────────


# extension 만으로 source_type 결정 — PDF 는 기존 ingest_pdf 와 일관성 위해 "pdf",
# 그 외 office/text 는 "document".
def _resolve_source_type(fmt: str, override: str | None) -> str:
    if override:
        return override
    if fmt == "pdf":
        return "pdf"
    return "document"


async def ingest_document(
    file_path: str | Path,
    *,
    source_type: str | None = None,         # caller 가 명시 안 하면 fmt 로 결정
    source_id: str | None = None,           # 채널별 외부 id (예: telegram msg_id)
    source_url: str | None = None,
    filename: str | None = None,            # 원본 파일명 (예: "photo_2026-05-18.jpg")
                                            # 없으면 file_path 의 basename 사용
    source_metadata_extra: dict[str, Any] | None = None,
    caption: str | None = None,              # 함께 온 사용자 메모 → items.user_notes
    analyze_now: bool = True,
    force: bool = False,
) -> dict[str, Any]:
    """다양한 포맷 file 을 LinkMind item 으로 ingest.

    사용 위치 (Phase 2.5+):
    - 텔레그램 첨부 (caption = 메시지 텍스트 → user_notes)
    - 향후: Slack/Discord 첨부, web upload, ...

    동작:
    1. 파일 bytes 로드 + storage 저장 (loss-less, sha256 dedup)
    2. extract_text_from_bytes 로 텍스트 추출
       - 지원 포맷이면 raw_content + chunks + summary
       - 지원 안 함 (image/doc/ppt/hwp/zip 등) → raw_content = 메타 placeholder,
         attachment 만 저장 (raw 파일은 보존 — 학습 데이터 손실 방지)
    3. items + attachments insert, dedup 시 attachment 만 추가
    4. caption 있으면 items.user_notes 에 단순 저장 (LLM 키워드 추출은 PATCH 시점에)
    5. analyze_now 면 chunks (embedding) + summary, PDF 면 figure 도

    반환:
        {item_id, created, source_type, source_format, file_hash, file_path,
         chunks_indexed, figures_saved, summary_generated, refreshed, ...}
    """
    # ── 1. 파일 로드 + storage 저장 ─────────────────────────────
    from backend.storage.local import save_file

    src = Path(file_path)
    if not src.exists():
        raise FileNotFoundError(f"file not found: {src}")

    storage_path, file_hash, file_size = save_file(src)
    data = src.read_bytes()        # 추출용 (storage 는 이미 복사됨)

    # ── 2. 포맷 식별 + 텍스트 추출 ──────────────────────────────
    # caller 가 원본 파일명 (예: 텔레그램 첨부의 "photo_2026-05-18.jpg") 을 명시했으면 우선,
    # 없으면 file_path 의 basename. 확장자 + title fallback + mime 추정 에 사용.
    effective_filename = filename or src.name
    if not _explicit_mime(source_metadata_extra):
        mime_type = mimetypes.guess_type(effective_filename)[0]
    else:
        mime_type = source_metadata_extra.get("mime_type")  # type: ignore[union-attr]

    fmt = guess_format(effective_filename, mime_type)
    extract = extract_text_from_bytes(data, filename=effective_filename, mime_type=mime_type)

    if extract:
        raw_content = extract.raw_text
        title = extract.title
        extractor = extract.extractor
        extract_meta = extract.meta
    else:
        # 텍스트 추출 안 됨 (image/unsupported) — raw 는 placeholder, 파일은 attachment 로 보존
        raw_content = (
            f"[binary file: {effective_filename}, mime={mime_type or 'unknown'}, "
            f"size={file_size} bytes, format={fmt}]"
        )
        title = _filename_to_title(effective_filename)
        extractor = "none"
        extract_meta = {}

    resolved_source_type = _resolve_source_type(fmt, source_type)
    content_hash = sha256_text(raw_content)

    base_metadata: dict[str, Any] = {
        "file_hash": file_hash,
        "file_size": file_size,
        "file_path": storage_path,
        "filename": effective_filename,
        "mime_type": mime_type,
        "source_format": fmt,
        "extractor": extractor,
        "extract_meta": extract_meta,
    }
    if source_metadata_extra:
        # caller (예: telegram) 가 channel/msg_id/caption 등 provenance 메타 추가.
        # 같은 key 충돌 시 caller 가 우선 (raw 원본 정보 보존).
        merged_meta = {**base_metadata, **source_metadata_extra}
    else:
        merged_meta = base_metadata

    final_source_url = source_url or f"/files/{file_hash}"

    # ── 3. DB insert / dedup ───────────────────────────────────
    engine = get_engine()
    session_factory = async_sessionmaker(engine, expire_on_commit=False, class_=AsyncSession)
    async with session_factory() as session:
        existing = await find_item_by_hash(
            session, source_type=resolved_source_type, content_hash=content_hash,
        )
        if existing is not None:
            # 같은 파일을 다른 메시지로 다시 보냈을 수도 — attachment dedup 자동 (UNIQUE).
            att_id = await insert_attachment(
                session, item_id=existing, file_path=storage_path,
                file_hash=file_hash, file_size=file_size,
                mime_type=mime_type or "application/octet-stream",
                role="attachment",
            )
            # 새 caption 이면 user_notes 에 append (기존 메모 보존, idempotent —
            # 같은 caption 두 번 들어와도 dedup. Phase 2.5 wave-3 정책)
            if caption and caption.strip():
                from backend.db.repository import append_item_user_notes
                await append_item_user_notes(
                    session, item_id=existing, new_note=caption.strip(),
                )
            await session.commit()
            return {
                "item_id": str(existing),
                "created": False,
                "source_type": resolved_source_type,
                "source_format": fmt,
                "file_hash": file_hash,
                "file_path": storage_path,
                "attachment_added": att_id is not None,
                "filename": filename,
                "chunks_indexed": 0,
                "figures_saved": 0,
            }

        # 본문에서 외부 ID (arxiv:xxx, doi:..., github:owner/repo, youtube id 등) 추출
        # — graph topic auto-link 의 단서. PDF 논문 본문에 arxiv id 있으면 같은 arxiv
        # topic 의 다른 modality (paper URL, GitHub repo) 와 자동으로 cluster.
        ext_ids = extract_external_ids(text=raw_content[:20000]) if raw_content else []

        merged_meta_with_ids = {
            **merged_meta,
            "external_ids": [{"kind": x.kind, "value": x.value} for x in ext_ids],
        }

        item_id = await insert_item(
            session,
            source_type=resolved_source_type,
            raw_content=raw_content,
            raw_content_hash=content_hash,
            source_id=source_id or file_hash,
            source_url=final_source_url,
            source_metadata=merged_meta_with_ids,
            title=title,
            source_created_at=None,
        )
        await insert_attachment(
            session, item_id=item_id, file_path=storage_path,
            file_hash=file_hash, file_size=file_size,
            mime_type=mime_type or "application/octet-stream",
            role="attachment",
        )
        # graph topic 자동 link (lazy import — backend.ingest.url 순환 의존 회피)
        from backend.ingest.url import auto_link_topics

        await auto_link_topics(
            session, item_id=item_id, source_type=resolved_source_type,
            title=title, ids=ext_ids,
        )
        # caption (텔레그램 등에서 함께 온 사용자 메모) → user_notes 자동.
        # 새 item 이라 기존 메모 없음 — append 도 같은 결과지만 일관성 위해 동일 API.
        if caption and caption.strip():
            from backend.db.repository import append_item_user_notes
            await append_item_user_notes(
                session, item_id=item_id, new_note=caption.strip(),
            )
        await session.commit()

        chunks_indexed = 0
        figures_saved = 0
        summary_generated = False
        if analyze_now and extract and extract.raw_text:
            # lazy import — 순환 의존 회피 (url 이 document import 할 수도 있어서)
            from backend.ingest.url import (
                ExtractedDoc,
                _embed_and_index,
                _generate_and_save_summary,
            )

            chunks_indexed = await _embed_and_index(
                session, item_id=item_id, text=raw_content,
            )
            if fmt == "pdf":
                figures_saved = await _save_pdf_figures(
                    session, item_id=item_id, data=data,
                )
                await session.commit()
            paper_keywords = [fmt]
            doc = ExtractedDoc(
                body=raw_content, title=title, abstract=None,
                paper_keywords=paper_keywords,
            )
            summary_text, _tags = await _generate_and_save_summary(
                session, item_id=item_id, doc=doc,
            )
            summary_generated = summary_text is not None

        return {
            "item_id": str(item_id),
            "created": True,
            "source_type": resolved_source_type,
            "source_format": fmt,
            "file_hash": file_hash,
            "file_path": storage_path,
            "filename": filename,
            "chunks_indexed": chunks_indexed,
            "figures_saved": figures_saved,
            "summary_generated": summary_generated,
            "title": title,
        }


def _explicit_mime(extra: dict[str, Any] | None) -> bool:
    """source_metadata_extra 에 mime_type 이 직접 명시됐는지 (telegram 등이 알면 우선)."""
    return bool(extra and extra.get("mime_type"))
