"""backend/ingest/document/__init__.py 단위 테스트.

format 식별 + 각 포맷의 in-memory fixture 로 텍스트 추출 검증. DOCX/PPTX 는
python-docx / python-pptx 가 직접 만들어주는 빈 문서 — 의존성 확인 + 추출 흐름
sanity. PDF 는 backend.ingest.pdf 의 fixture 가 별도 (tests/test_pdf_*) — 여기는
re-export 흐름만 mocking.
"""

from __future__ import annotations

import io

import pytest

from backend.ingest.document import (
    DocumentExtract,
    extract_text_from_bytes,
    guess_format,
    is_supported,
)


# ── guess_format ─────────────────────────────────────────────


def test_guess_format_by_mime():
    assert guess_format(mime_type="application/pdf") == "pdf"
    assert guess_format(
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ) == "docx"
    assert guess_format(
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ) == "pptx"
    assert guess_format(mime_type="text/plain") == "txt"
    assert guess_format(mime_type="text/markdown") == "markdown"
    assert guess_format(mime_type="image/png") == "image"


def test_guess_format_by_extension():
    assert guess_format(filename="paper.pdf") == "pdf"
    assert guess_format(filename="report.docx") == "docx"
    assert guess_format(filename="slides.pptx") == "pptx"
    assert guess_format(filename="memo.txt") == "txt"
    assert guess_format(filename="README.md") == "markdown"
    assert guess_format(filename="photo.JPG") == "image"


def test_guess_format_mime_priority_over_ext():
    """mime_type 이 우선 — 파일명 확장자 잘못된 경우 보호."""
    assert guess_format(filename="thing.docx", mime_type="application/pdf") == "pdf"


def test_guess_format_old_office_marked_separately():
    """구 binary 형식은 doc/ppt 로 — Phase 3+ 에 LibreOffice 처리."""
    assert guess_format(filename="legacy.doc") == "doc"
    assert guess_format(filename="legacy.ppt") == "ppt"
    assert guess_format(filename="report.hwp") == "hwp"
    assert guess_format(filename="report.hwpx") == "hwp"


def test_guess_format_unknown_returns_unknown():
    assert guess_format(filename="archive.zip") == "unknown"
    assert guess_format() == "unknown"


def test_is_supported_set():
    assert is_supported("pdf") is True
    assert is_supported("docx") is True
    assert is_supported("pptx") is True
    assert is_supported("txt") is True
    assert is_supported("markdown") is True

    assert is_supported("doc") is False
    assert is_supported("ppt") is False
    assert is_supported("hwp") is False
    assert is_supported("image") is False
    assert is_supported("unknown") is False


# ── extract_text_from_bytes — text formats ───────────────────


def test_extract_txt_utf8():
    data = "안녕하세요\n포인트클라우드 메모".encode("utf-8")
    res = extract_text_from_bytes(data, filename="memo.txt")
    assert isinstance(res, DocumentExtract)
    assert "포인트클라우드" in res.raw_text
    assert res.source_format == "txt"
    assert res.title == "memo"


def test_extract_txt_cp949_fallback():
    """cp949 인코딩 (한국 윈도우 메모장 기본). utf-8 실패 후 charset-normalizer."""
    pytest.importorskip("charset_normalizer")
    data = "한국어 텍스트 메모".encode("cp949")
    res = extract_text_from_bytes(data, filename="windows_memo.txt")
    assert isinstance(res, DocumentExtract)
    assert "한국어" in res.raw_text


def test_extract_markdown_extracts_first_h1_as_title():
    data = b"# My Paper Title\n\nbody content here"
    res = extract_text_from_bytes(data, filename="paper.md")
    assert isinstance(res, DocumentExtract)
    assert res.source_format == "markdown"
    assert res.title == "My Paper Title"
    assert "body content" in res.raw_text


def test_extract_markdown_no_h1_falls_back_to_filename():
    data = b"plain markdown without heading"
    res = extract_text_from_bytes(data, filename="readme.md")
    assert isinstance(res, DocumentExtract)
    assert res.title == "readme"


# ── extract_text_from_bytes — docx / pptx (in-memory fixture) ─


def _make_docx_bytes(*paragraphs: str, title: str | None = None) -> bytes:
    """python-docx 로 메모리 안에서 .docx 만들어 bytes 로 반환."""
    from docx import Document
    doc = Document()
    if title:
        doc.core_properties.title = title
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_extract_docx_basic():
    data = _make_docx_bytes(
        "포인트클라우드 압축 연구",
        "Section 1: introduction",
        title="My Research Note",
    )
    res = extract_text_from_bytes(data, filename="note.docx")
    assert isinstance(res, DocumentExtract)
    assert res.source_format == "docx"
    assert res.extractor == "python-docx"
    assert "포인트클라우드" in res.raw_text
    assert "introduction" in res.raw_text
    assert res.title == "My Research Note"


def test_extract_docx_empty_returns_none():
    """본문 비어있으면 None — caller 가 attachment 만 저장."""
    data = _make_docx_bytes()
    res = extract_text_from_bytes(data, filename="empty.docx")
    assert res is None


def _make_pptx_bytes(*slides: list[str]) -> bytes:
    """python-pptx 로 메모리 안에서 .pptx 만듦. 각 slide 는 텍스트 list."""
    from pptx import Presentation
    pres = Presentation()
    blank_layout = pres.slide_layouts[6]  # blank
    for texts in slides:
        slide = pres.slides.add_slide(blank_layout)
        for t in texts:
            from pptx.util import Inches
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tb.text_frame.text = t
    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


def test_extract_pptx_basic():
    data = _make_pptx_bytes(
        ["3D Gaussian Splatting", "Title slide"],
        ["Method overview", "어떻게 동작하나"],
    )
    res = extract_text_from_bytes(data, filename="slides.pptx")
    assert isinstance(res, DocumentExtract)
    assert res.source_format == "pptx"
    assert res.extractor == "python-pptx"
    assert "3D Gaussian Splatting" in res.raw_text
    assert "어떻게 동작" in res.raw_text
    # 첫 슬라이드 첫 텍스트가 title
    assert res.title == "3D Gaussian Splatting"
    assert res.meta.get("slide_count") == 2


# ── extract_text_from_bytes — unsupported / edge ─────────────


def test_extract_unknown_format_returns_none():
    """지원 안 하는 포맷은 None — caller 가 attachment 만 저장."""
    res = extract_text_from_bytes(b"binary data", filename="archive.zip")
    assert res is None


def test_extract_image_returns_none():
    """이미지는 None (Phase 3+ OCR)."""
    res = extract_text_from_bytes(b"fake-png-bytes", filename="photo.png")
    assert res is None


def test_extract_empty_bytes_returns_none():
    assert extract_text_from_bytes(b"", filename="empty.pdf") is None


def test_extract_broken_docx_returns_none():
    """깨진 DOCX 도 graceful — None 반환, 예외 X."""
    res = extract_text_from_bytes(b"not a real docx", filename="broken.docx")
    assert res is None
